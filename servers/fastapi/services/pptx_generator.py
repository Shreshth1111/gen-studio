"""
═══════════════════════════════════════════════════════════════════════════════
 pptx_generator.py — EXPORT: SLIDE JSON  →  REAL .pptx FILE
═══════════════════════════════════════════════════════════════════════════════

WHAT THIS FILE DOES
  Takes the same slide JSON the browser renders and rebuilds it as a NATIVE
  PowerPoint file using python-pptx, so the downloaded .pptx looks like the
  on-screen deck and is fully editable in PowerPoint/Keynote/Google Slides.

WHERE IT FITS
  api/v1/ppt/endpoints/export.py  →  POST /export/{id}/pptx
     builds ppt_data {title, theme, slides:[...]} from the DB and calls
     PPTXGenerator(theme).generate(ppt_data, output_path). The .pptx is written
     under app_data/presentations/ and served for download (and optionally
     handed to SageStudio for video lectures).

HOW IT WORKS
  • Canvas: a blank 16:9 slide (13.33in × 7.5in), one per slide_data.
  • generate() dispatches on slide_data["layout_type"] to a `_render_<layout>`
    method — one renderer per layout, mirroring the frontend SlideRenderer.tsx
    switch, so the same 30+ layouts exist on BOTH sides.
  • Shared drawing helpers:
       _add_text_box / _add_bullets / _add_runs_box  — text (with rich runs)
       _add_image      — resolves /app_data/images/<id>.png → disk path
       _add_card / _add_accent_bar / _add_filled_shape / _add_polygon — shapes
       _theme_chart / _add_native_chart / _add_pie_chart — native PPTX charts
       _rgb(key) / _contrast_on — pull/choose colours from the theme palette
  • Colours come entirely from the `theme` dict (see lib/themes.py).

TECH: python-pptx (Inches/Pt/Emu units, RGBColor, MSO_SHAPE, CategoryChartData).
═══════════════════════════════════════════════════════════════════════════════
"""
import os
from pathlib import Path
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from typing import Optional

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


def _resolve_local_image(image_path_or_url: str) -> Optional[str]:
    """Map an image_url like /app_data/images/<id>.png to its on-disk path."""
    if not image_path_or_url:
        return None
    if image_path_or_url.startswith("http"):
        return None
    app_data = Path(os.getenv("APP_DATA_DIR", "./app_data")).resolve()
    if image_path_or_url.startswith("/app_data/"):
        rel = image_path_or_url[len("/app_data/"):]
        return str(app_data / rel)
    if image_path_or_url.startswith("/"):
        # absolute filesystem path
        return image_path_or_url
    return image_path_or_url


class PPTXGenerator:
    def __init__(self, theme: dict):
        self.theme = theme

    def generate(self, ppt_data: dict, output_path: str) -> str:
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        blank_layout = prs.slide_layouts[6]

        # Paginate at export time: code/bullet slides that are too long to fit
        # one printed slide are split into continuation slides ("(1/2)", "(2/2)").
        # This works for EVERY deck on download — even ones generated before the
        # split logic existed — since the .pptx can't scroll like the browser.
        for slide_data in self._paginate_for_export(ppt_data.get("slides", [])):
            slide = prs.slides.add_slide(blank_layout)
            self._apply_background(slide, slide_data)
            layout = slide_data.get("layout_type", "bullets")
            handler = getattr(self, f"_render_{layout}", self._render_bullets)
            handler(slide, slide_data)
            # Free-positioned image overlays the user placed on top of the
            # layout. Coordinates are normalised 0-1; multiply by slide size.
            self._render_overlays(slide, slide_data)
            if slide_data.get("speaker_notes"):
                slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

        prs.save(output_path)
        return output_path

    @staticmethod
    def _paginate_for_export(slides: list) -> list:
        """Split overflowing `code` (>22 lines) and `bullets` (>7) slides into
        continuation slides so nothing runs off the page in PowerPoint."""
        import re as _re
        MAX_CODE_LINES = 14
        MAX_BULLETS = 7
        out: list = []
        for s in slides:
            content = s.get("content") or {}
            layout = s.get("layout_type")
            code = content.get("code")
            bullets = content.get("bullets")

            chunks = None
            key = join = None
            if layout == "code" and isinstance(code, str) and code.count("\n") + 1 > MAX_CODE_LINES:
                lines = code.split("\n")
                chunks = [lines[i:i + MAX_CODE_LINES] for i in range(0, len(lines), MAX_CODE_LINES)]
                key, join = "code", "\n"
            elif layout == "bullets" and isinstance(bullets, list) and len(bullets) > MAX_BULLETS:
                chunks = [bullets[i:i + MAX_BULLETS] for i in range(0, len(bullets), MAX_BULLETS)]
                key, join = "bullets", None

            if not chunks:
                out.append(s)
                continue

            total = len(chunks)
            # Strip an existing "(n/m)" suffix so re-paginating an already-split
            # slide doesn't stack suffixes ("Foo (1/2) (1/2)").
            base_title = (content.get("title") or s.get("title") or "Slide")
            base_title = _re.sub(r"\s*\(\d+/\d+\)\s*$", "", base_title)
            for idx, chunk in enumerate(chunks):
                part = dict(s)
                new_content = dict(content)
                new_content[key] = join.join(chunk) if join is not None else chunk
                title = base_title if total == 1 else f"{base_title} ({idx + 1}/{total})"
                new_content["title"] = title
                part["title"] = title
                part["content"] = new_content
                if idx > 0:
                    part["image_url"] = None
                out.append(part)
        return out

    def _render_overlays(self, slide, slide_data):
        """Place free-positioned image overlays from ``content.overlays`` on
        top of the layout. Each overlay is ``{type:'image', src, x, y, w, h}``
        with normalised 0-1 coordinates."""
        content = slide_data.get("content", {}) or {}
        overlays = content.get("overlays") or []
        if not isinstance(overlays, list):
            return
        SLIDE_W_IN = 13.33
        SLIDE_H_IN = 7.5
        ordered = sorted(overlays, key=lambda o: (o or {}).get("z", 0))
        for o in ordered:
            if not isinstance(o, dict):
                continue
            if o.get("type") not in (None, "image"):
                continue
            src = o.get("src")
            if not src:
                continue
            try:
                x = float(o.get("x", 0.1))
                y = float(o.get("y", 0.1))
                w = float(o.get("w", 0.3))
                h = float(o.get("h", 0.3))
            except (TypeError, ValueError):
                continue
            self._add_image(
                slide, src,
                max(0.0, x) * SLIDE_W_IN,
                max(0.0, y) * SLIDE_H_IN,
                max(0.05, w) * SLIDE_W_IN,
                max(0.05, h) * SLIDE_H_IN,
            )

    def _apply_background(self, slide, slide_data):
        # Themes that define a radial gradient (royal/ocean/corporate_red) get
        # a real gradient fill so the .pptx matches the browser. Others stay
        # solid. Any failure falls back to a solid bg so export never breaks.
        grad = self.theme.get("bgGradient")
        if grad:
            try:
                self._apply_gradient_background(slide, grad)
                return
            except Exception:
                pass
        bg_color = self.theme.get("bg", "#FFFFFF").lstrip("#")
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(bg_color)

    def _apply_gradient_background(self, slide, stops):
        """Set a radial (path) gradient on the slide background by writing the
        OOXML ``<a:gradFill>`` directly. ``stops`` is a list of (hex, pos%)."""
        from pptx.oxml.ns import qn
        from lxml import etree

        fill = slide.background.fill
        fill.gradient()                      # ensure a <a:gradFill> exists
        bgPr = fill._xPr                      # the <p:bgPr> element
        old = bgPr.find(qn("a:gradFill"))
        if old is not None:
            bgPr.remove(old)

        A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        grad = etree.SubElement(bgPr, qn("a:gradFill"))
        gsLst = etree.SubElement(grad, qn("a:gsLst"))
        for hexc, pos in stops:
            gs = etree.SubElement(gsLst, qn("a:gs"))
            gs.set("pos", str(int(pos) * 1000))          # 0–100 → 0–100000
            clr = etree.SubElement(gs, qn("a:srgbClr"))
            clr.set("val", hexc.lstrip("#").upper())
        # Radial gradient with the focus near the top-left (≈30%/20%),
        # mirroring the CSS `radial-gradient(circle at 30% 20%, …)`.
        path = etree.SubElement(grad, qn("a:path"))
        path.set("path", "circle")
        rect = etree.SubElement(path, qn("a:fillToRect"))
        rect.set("l", "30000"); rect.set("t", "20000")
        rect.set("r", "30000"); rect.set("b", "20000")
        # gradFill must sit before any <a:effectLst>; bgPr children order is fine
        # as inserted (gradFill first). PowerPoint tolerates this ordering.

    def _rgb(self, key: str) -> RGBColor:
        return RGBColor.from_string(self.theme.get(key, "#000000").lstrip("#"))

    @staticmethod
    def _contrast_on(hex_color: str) -> str:
        """Return '#111111' or '#FFFFFF' — whichever reads better on top of
        `hex_color`. Mirrors the frontend `contrastOn` so funnel / shape
        labels are legible and identical across browser and PPTX."""
        h = (hex_color or "#000000").lstrip("#")
        if len(h) < 6:
            return "#FFFFFF"
        r = int(h[0:2], 16); g = int(h[2:4], 16); b = int(h[4:6], 16)
        lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        return "#111111" if lum > 0.6 else "#FFFFFF"

    def _add_text_box(self, slide, text, left, top, width, height,
                      font_size=18, bold=False, align=PP_ALIGN.LEFT, color_key="text",
                      color_hex=None):
        """Add a text frame.

        Now rich-text aware: if ``text`` contains HTML tags, it is parsed into
        paragraphs of styled runs (bold/italic/underline/strikethrough/colour/
        link/alignment + bulleted lists) using ``utils.html_runs``. Plain
        strings still work exactly as before — they're split on ``\\n``.
        """
        from utils.html_runs import parse as parse_runs
        from pptx.dml.color import RGBColor

        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
        except Exception:
            pass
        try:
            from pptx.util import Emu
            tf.margin_left   = Emu(36000)
            tf.margin_right  = Emu(36000)
            tf.margin_top    = Emu(18000)
            tf.margin_bottom = Emu(18000)
        except Exception:
            pass

        paragraphs = parse_runs(text or "")
        if not paragraphs:
            # Empty input — still emit one empty paragraph so the box exists.
            paragraphs = [type("P", (), {"runs": [type("R", (), {
                "text": "", "bold": False, "italic": False, "underline": False,
                "strike": False, "color": None, "link": None,
            })()], "align": None, "is_list_item": False})()]

        # Default colour used when a run carries no explicit colour. An
        # explicit color_hex overrides the theme color_key (used by funnel
        # slice labels that need a contrast-computed colour).
        if color_hex:
            default_rgb = RGBColor.from_string(color_hex.lstrip("#"))
        else:
            default_rgb = self._rgb(color_key)

        for i, para in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Pick alignment: explicit in the HTML wins, otherwise caller default.
            if para.align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif para.align == "right":
                p.alignment = PP_ALIGN.RIGHT
            elif para.align == "left":
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = align

            # Clear any auto-created blank run so we don't bake formatting holes.
            for run in list(p.runs):
                run.text = ""

            for r in para.runs:
                if not r.text:
                    continue
                run = p.add_run()
                run.text = r.text
                run.font.size = Pt(font_size)
                # Caller-default bold OR per-run bold — caller wins as the
                # baseline weight, per-run can add emphasis.
                run.font.bold = bool(bold or r.bold)
                if r.italic:
                    run.font.italic = True
                if r.underline:
                    run.font.underline = True
                # python-pptx doesn't expose strikethrough on `Font` directly;
                # patch it via the raw XML attribute when needed.
                if r.strike:
                    try:
                        from pptx.oxml.ns import qn
                        rPr = run.font._rPr
                        if rPr is not None:
                            rPr.set("strike", "sngStrike")
                    except Exception:
                        pass
                # Colour: explicit run colour wins; otherwise the caller's
                # color_key (theme) provides the baseline.
                if r.color and r.color.startswith("#"):
                    from pptx.dml.color import RGBColor
                    try:
                        run.font.color.rgb = RGBColor.from_string(r.color.lstrip("#"))
                    except Exception:
                        run.font.color.rgb = default_rgb
                else:
                    run.font.color.rgb = default_rgb
                # Hyperlink — uses python-pptx's run.hyperlink helper.
                if r.link:
                    try:
                        run.hyperlink.address = r.link
                    except Exception:
                        pass
        return tf

    def _add_bullets(self, slide, bullets, left, top, width, height,
                     font_size=18, bullet_key="accent", text_key="text"):
        """Render a bulleted list where the bullet glyph is painted in the
        accent colour and the text in the body colour — matching the browser,
        which draws an accent dot beside each bullet. Each bullet string may
        itself contain rich-text HTML, which is parsed and styled per-run."""
        from utils.html_runs import parse as parse_runs
        from pptx.dml.color import RGBColor

        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
        except Exception:
            pass
        try:
            from pptx.util import Emu
            tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
            tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
        except Exception:
            pass

        bullet_rgb = self._rgb(bullet_key)
        default_rgb = self._rgb(text_key)
        first = True
        for b in (bullets or []):
            text = b if isinstance(b, str) else str(b)
            # Flatten the bullet's (possibly multi-paragraph) HTML into one line.
            runs = [r for para in parse_runs(text) for r in para.runs]
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            for run in list(p.runs):
                run.text = ""
            # Accent-coloured bullet glyph.
            brun = p.add_run()
            brun.text = "●  "
            brun.font.size = Pt(font_size)
            brun.font.color.rgb = bullet_rgb
            # Body runs (rich text), defaulting to the body colour.
            for r in runs:
                if not r.text:
                    continue
                run = p.add_run()
                run.text = r.text
                run.font.size = Pt(font_size)
                run.font.bold = bool(r.bold)
                if r.italic:
                    run.font.italic = True
                if r.underline:
                    run.font.underline = True
                if r.strike:
                    try:
                        rPr = run.font._rPr
                        if rPr is not None:
                            rPr.set("strike", "sngStrike")
                    except Exception:
                        pass
                if r.color and r.color.startswith("#"):
                    try:
                        run.font.color.rgb = RGBColor.from_string(r.color.lstrip("#"))
                    except Exception:
                        run.font.color.rgb = default_rgb
                else:
                    run.font.color.rgb = default_rgb
                if r.link:
                    try:
                        run.hyperlink.address = r.link
                    except Exception:
                        pass
        return tf

    def _add_runs_box(self, slide, runs_per_line, left, top, width, height,
                      font_size=14, align=PP_ALIGN.LEFT):
        """Like `_add_text_box` but each *line* is a list of (text, color_hex,
        bold) runs. Used when one paragraph needs differently-coloured pieces
        — e.g. comparison rows with a green ✓ marker followed by body text."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE
        except Exception:
            pass
        try:
            tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
            tf.margin_top  = Emu(18000); tf.margin_bottom = Emu(18000)
        except Exception:
            pass
        for i, runs in enumerate(runs_per_line):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            for run in list(p.runs):
                run.text = ""
            for text, color_hex, bold in runs:
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = RGBColor.from_string(color_hex.lstrip("#"))
        return tf

    def _add_image(self, slide, image_path_or_url, left, top, width, height):
        if not image_path_or_url:
            return
        try:
            if image_path_or_url.startswith("http"):
                import urllib.request
                local_path, _ = urllib.request.urlretrieve(image_path_or_url)
            else:
                local_path = _resolve_local_image(image_path_or_url) or image_path_or_url
            if not local_path or not Path(local_path).exists():
                print(f"Image not found on disk: {image_path_or_url} -> {local_path}")
                return
            slide.shapes.add_picture(
                local_path, Inches(left), Inches(top), Inches(width), Inches(height)
            )
        except Exception as e:
            print(f"Image add failed: {e}")

    def _add_accent_bar(self, slide, left, top, width, height, color_key="accent"):
        """Decorative thin accent bar — used as a section divider in slides."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(color_key)
        return shape

    def _add_card(self, slide, left, top, width, height, color_key="secondary", radius=True):
        """Card-style rounded rectangle used as a background for stats/columns."""
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(color_key)
        return shape

    def _render_title(self, slide, slide_data):
        content = slide_data.get("content", {})
        # Decorative left accent strip
        self._add_accent_bar(slide, 0.0, 0.0, 0.4, 7.5, color_key="accent")
        title = content.get("title", slide_data.get("title", ""))
        self._add_text_box(slide, title,
                           1.0, 2.5, 11.33, 1.6, font_size=54, bold=True,
                           align=PP_ALIGN.CENTER, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 1.5, 4.3, 10.33, 1.0,
                               font_size=22, align=PP_ALIGN.CENTER, color_key="muted")
        # Bottom underline
        self._add_accent_bar(slide, 6.16, 5.6, 1.0, 0.07, color_key="accent")

    def _render_bullets(self, slide, slide_data):
        content = slide_data.get("content", {})
        title = content.get("title", slide_data.get("title", ""))
        # Title with accent underline
        self._add_text_box(slide, title, 0.6, 0.45, 12.33, 0.9,
                           font_size=34, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.25, 0.7, 0.06, color_key="accent")

        bullets = content.get("bullets", [])
        has_image = bool(content.get("image_url"))
        text_width = 7.0 if has_image else 12.0
        # Pretty bulleted text — Pt sizes adapt to count
        font_size = 22 if len(bullets) <= 4 else 18
        self._add_bullets(slide, bullets, 0.6, 1.55, text_width, 5.5,
                          font_size=font_size)
        if has_image:
            self._add_image(slide, content["image_url"], 7.9, 1.55, 5.0, 5.3)
        if content.get("callout"):
            # Match the on-screen callout: secondary bg + accent left edge +
            # heading-coloured copy.
            self._add_card(slide, 0.6, 6.5, text_width, 0.6, color_key="secondary")
            self._add_accent_bar(slide, 0.6, 6.5, 0.08, 0.6, color_key="accent")
            self._add_text_box(slide, content["callout"], 0.85, 6.6, text_width-0.3, 0.5,
                               font_size=14, bold=True, color_key="heading")

    def _render_two_column(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")

        # Two cards side by side
        for i, key in enumerate(["col1", "col2"]):
            left = 0.6 + i * 6.15
            self._add_card(slide, left, 1.5, 6.0, 5.6, color_key="secondary")
            self._add_text_box(slide, content.get(f"{key}_heading", ""),
                               left + 0.3, 1.7, 5.6, 0.6,
                               font_size=20, bold=True, color_key="accent")
            bullets = content.get(f"{key}_bullets", [])
            self._add_bullets(slide, bullets, left + 0.3, 2.4, 5.6, 4.6,
                              font_size=15)

    def _render_image_left(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        if content.get("image_url"):
            self._add_image(slide, content["image_url"], 0.6, 1.5, 5.6, 5.6)
        else:
            self._add_card(slide, 0.6, 1.5, 5.6, 5.6, color_key="secondary")
        if content.get("content_heading"):
            self._add_text_box(slide, content["content_heading"], 6.4, 1.6, 6.4, 0.6,
                               font_size=22, bold=True, color_key="accent")
        self._add_bullets(slide, content.get("bullets", []), 6.4, 2.4, 6.4, 4.7, font_size=18)

    def _render_image_right(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        if content.get("content_heading"):
            self._add_text_box(slide, content["content_heading"], 0.6, 1.6, 6.4, 0.6,
                               font_size=22, bold=True, color_key="accent")
        self._add_bullets(slide, content.get("bullets", []), 0.6, 2.4, 6.4, 4.7, font_size=18)
        if content.get("image_url"):
            self._add_image(slide, content["image_url"], 7.2, 1.5, 5.6, 5.6)
        else:
            self._add_card(slide, 7.2, 1.5, 5.6, 5.6, color_key="secondary")

    def _render_stats(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=32, bold=True, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 0.6, 1.15, 12.33, 0.5,
                               font_size=16, color_key="muted")
        self._add_accent_bar(slide, 0.6, 1.7, 0.7, 0.06, color_key="accent")
        stats = content.get("stats", [])
        n = max(min(len(stats), 4), 1)
        gap = 0.25
        total_w = 12.33 - gap * (n - 1)
        col_width = total_w / n
        for i, stat in enumerate(stats[:n]):
            left = 0.5 + i * (col_width + gap)
            self._add_card(slide, left, 2.2, col_width, 4.4, color_key="secondary")
            self._add_text_box(slide, stat.get("value", ""), left, 2.7, col_width, 1.8,
                               font_size=64, bold=True, align=PP_ALIGN.CENTER, color_key="accent")
            self._add_text_box(slide, stat.get("label", ""), left, 4.6, col_width, 0.6,
                               font_size=18, bold=True, align=PP_ALIGN.CENTER, color_key="heading")
            self._add_text_box(slide, stat.get("context", ""), left + 0.2, 5.3, col_width - 0.4, 1.0,
                               font_size=12, align=PP_ALIGN.CENTER, color_key="muted")

    def _render_quote(self, slide, slide_data):
        content = slide_data.get("content", {})
        # Big decorative quote mark
        self._add_text_box(slide, "“", 0.8, 0.8, 2.0, 2.0,
                           font_size=140, bold=True, color_key="accent")
        self._add_text_box(slide, content.get("quote", ""),
                           1.5, 2.3, 10.33, 3.0, font_size=30, bold=True,
                           align=PP_ALIGN.LEFT, color_key="heading")
        self._add_accent_bar(slide, 1.5, 5.5, 0.6, 0.06, color_key="accent")
        attribution = content.get("attribution", "")
        if attribution:
            self._add_text_box(slide, attribution, 1.5, 5.7, 10.33, 0.5,
                               font_size=18, bold=True, color_key="accent")
        if content.get("role"):
            self._add_text_box(slide, content["role"], 1.5, 6.2, 10.33, 0.5,
                               font_size=14, color_key="muted")

    def _render_timeline(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        events = content.get("events", [])[:6]
        if not events:
            return
        # Match the on-screen renderer: dot ABOVE year/label/description, with
        # an accent-coloured ring whose centre is a small `bg`-coloured dot.
        dot_y = 2.0
        dot_size = 0.42
        inner_size = 0.18
        # Horizontal connector line passes through the centre of the dots.
        line_y = dot_y + dot_size / 2 - 0.02
        self._add_accent_bar(slide, 0.6, line_y, 12.13, 0.04, color_key="border")
        col_w = 12.13 / len(events)
        for i, evt in enumerate(events):
            left = 0.6 + i * col_w
            cx = left + (col_w / 2) - dot_size / 2
            # Outer ring — accent
            outer = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(cx), Inches(dot_y),
                Inches(dot_size), Inches(dot_size)
            )
            outer.line.fill.background()
            outer.fill.solid()
            outer.fill.fore_color.rgb = self._rgb("accent")
            # Inner dot — bg
            inner_off = (dot_size - inner_size) / 2
            inner = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx + inner_off), Inches(dot_y + inner_off),
                Inches(inner_size), Inches(inner_size)
            )
            inner.line.fill.background()
            inner.fill.solid()
            inner.fill.fore_color.rgb = self._rgb("bg")
            # Year below the dot
            self._add_text_box(slide, evt.get("year", ""), left, dot_y + dot_size + 0.15,
                               col_w, 0.6,
                               font_size=20, bold=True, align=PP_ALIGN.CENTER, color_key="accent")
            # Label
            self._add_text_box(slide, evt.get("label", ""), left,
                               dot_y + dot_size + 0.85, col_w, 0.6,
                               font_size=15, bold=True,
                               align=PP_ALIGN.CENTER, color_key="heading")
            # Description
            self._add_text_box(slide, evt.get("description", ""),
                               left + 0.1, dot_y + dot_size + 1.55, col_w - 0.2, 2.5,
                               font_size=11, align=PP_ALIGN.CENTER, color_key="muted")

    def _render_comparison(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")

        # Match Tailwind's emerald-500 / red-500 for the per-row tick marks.
        EMERALD = "10B981"
        RED     = "EF4444"
        text_hex    = self.theme.get("text", "#1A1A1A").lstrip("#")

        for i, side in enumerate(["option_a", "option_b"]):
            left = 0.6 + i * 6.15
            opt = content.get(side, {})
            self._add_card(slide, left, 1.5, 6.0, 5.6, color_key="secondary")
            self._add_text_box(slide, opt.get("label", ""), left + 0.3, 1.7, 5.6, 0.6,
                               font_size=22, bold=True, color_key="accent")
            self._add_accent_bar(slide, left + 0.3, 2.4, 0.5, 0.04, color_key="accent")

            self._add_text_box(slide, "PROS", left + 0.3, 2.6, 5.6, 0.4,
                               font_size=11, bold=True, color_key="muted")
            pros_lines = [
                [("✓  ", EMERALD, True), (p, text_hex, False)]
                for p in opt.get("pros", [])
            ]
            if pros_lines:
                self._add_runs_box(slide, pros_lines,
                                   left + 0.3, 3.0, 5.6, 1.7, font_size=14)

            self._add_text_box(slide, "CONS", left + 0.3, 4.8, 5.6, 0.4,
                               font_size=11, bold=True, color_key="muted")
            cons_lines = [
                [("✗  ", RED, True), (c, text_hex, False)]
                for c in opt.get("cons", [])
            ]
            if cons_lines:
                self._add_runs_box(slide, cons_lines,
                                   left + 0.3, 5.2, 5.6, 1.7, font_size=14)

    def _render_team(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", "Meet the Team"), 0.6, 0.4, 12.33, 0.7,
                           font_size=32, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        members = content.get("members", [])
        n = max(min(len(members), 4), 1)
        gap = 0.25
        col_w = (12.33 - gap * (n - 1)) / n
        for i, m in enumerate(members[:n]):
            left = 0.5 + i * (col_w + gap)
            self._add_card(slide, left, 1.7, col_w, 5.0, color_key="secondary")
            # Avatar circle (initial)
            avatar_size = 1.4
            avatar_left = left + (col_w - avatar_size) / 2
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(avatar_left), Inches(2.0),
                Inches(avatar_size), Inches(avatar_size)
            )
            shape.line.fill.background()
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb("accent")
            initial = (m.get("name") or "?")[0].upper()
            self._add_text_box(slide, initial, avatar_left, 2.15, avatar_size, 1.1,
                               font_size=44, bold=True, align=PP_ALIGN.CENTER, color_key="bg")
            self._add_text_box(slide, m.get("name", ""), left, 3.6, col_w, 0.5,
                               font_size=18, bold=True, align=PP_ALIGN.CENTER, color_key="heading")
            self._add_text_box(slide, m.get("role", ""), left, 4.15, col_w, 0.4,
                               font_size=13, align=PP_ALIGN.CENTER, color_key="accent")
            self._add_text_box(slide, m.get("bio", ""), left + 0.2, 4.7, col_w - 0.4, 1.8,
                               font_size=11, align=PP_ALIGN.CENTER, color_key="muted")

    def _render_agenda(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", "Agenda"), 0.6, 0.4, 5.5, 1.5,
                           font_size=44, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.6, 0.7, 0.06, color_key="accent")
        items = content.get("items", [])[:8]
        n = max(len(items), 1)
        spacing = min(0.85, 5.0 / n)
        start_top = 2.0 if n <= 5 else 1.4
        for i, item in enumerate(items):
            top = start_top + i * spacing
            self._add_text_box(slide, item.get("number", str(i+1).zfill(2)),
                               6.0, top, 1.0, 0.7,
                               font_size=28, bold=True, color_key="accent")
            self._add_text_box(slide, item.get("label", ""), 7.0, top + 0.05, 6.0, 0.6,
                               font_size=20, color_key="text")

    def _render_blank(self, slide, slide_data):
        pass

    def _render_code(self, slide, slide_data):
        """Dedicated code slide: title + a dark monospace code panel that
        preserves indentation and newlines (never bullets)."""
        content = slide_data.get("content", {}) or {}
        title = content.get("title") or slide_data.get("title") or ""
        language = (content.get("language") or "code")
        code = content.get("code") or ""
        caption = content.get("caption") or ""

        # Header
        self._add_accent_bar(slide, 0.0, 0.0, 0.18, 7.5)
        self._add_text_box(slide, title, 0.7, 0.5, 11.9, 0.9,
                           font_size=30, bold=True, color_key="heading")
        # Language chip
        self._add_text_box(slide, language.upper(), 0.7, 1.35, 3.0, 0.4,
                           font_size=11, bold=True, color_key="accent")

        # Dark code panel
        panel_top, panel_h = 1.85, 4.7 if not caption else 4.3
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(panel_top),
            Inches(11.93), Inches(panel_h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor.from_string("0D1117")  # GitHub-dark
        panel.line.color.rgb = self._rgb("border")
        panel.line.width = Pt(1)
        panel.shadow.inherit = False

        # Code text — monospace, preserve every line. Unlike the browser (which
        # scrolls), a PPTX text box can't scroll, so the font is sized DOWN to
        # fit the panel both vertically (line count) and horizontally (longest
        # line). PowerPoint's shrink-to-fit is also enabled as a backstop.
        inner_top, inner_h = panel_top + 0.18, panel_h - 0.36
        inner_left, inner_w = 1.0, 11.3
        tb = slide.shapes.add_textbox(Inches(inner_left), Inches(inner_top),
                                      Inches(inner_w), Inches(inner_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(4)
        tf.margin_top = tf.margin_bottom = Pt(2)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink text on overflow
        except Exception:
            pass

        lines = code.split("\n") if code else [""]
        n_lines = max(1, len(lines))
        longest = max((len(ln) for ln in lines), default=1)
        line_spacing = 1.2
        # Font is sized CONSERVATIVELY so it fits even when PowerPoint
        # substitutes a slightly wider/taller monospace font and doesn't
        # re-run autofit (e.g. in Protected View). Long code is split into
        # ≤14-line chunks upstream, so this stays a readable ~12-13pt.
        #   • usable height = 88% of the panel (safety margin)
        #   • per-line budget uses 1.35× (bigger than the actual 1.2 spacing)
        safe_h = inner_h * 0.88
        size_by_h = (safe_h * 72.0) / (n_lines * 1.35)
        size_by_w = ((inner_w - 0.2) * 72.0) / (max(longest, 1) * 0.62)  # Consolas ≈ 0.62·size wide
        font_pt = max(8.0, min(13.0, size_by_h, size_by_w))

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line if line != "" else " "
            run.font.name = "Consolas"
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor.from_string("E6EDF3")
            p.line_spacing = line_spacing

        if caption:
            self._add_text_box(slide, caption, 0.7, panel_top + panel_h + 0.15, 11.9, 0.5,
                               font_size=13, color_key="muted")

    def _render_table(self, slide, slide_data):
        """Professional comparison table using a native PowerPoint table so it
        stays editable. Header row = accent fill + contrast text; first column
        emphasised; body rows zebra-striped to match the browser renderer.

        Schema: {title, headers:[...], rows:[[...],[...]]}.
        """
        from pptx.util import Emu, Pt as _Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN as _AL, MSO_ANCHOR

        content = slide_data.get("content", {})
        title = content.get("title", slide_data.get("title", ""))
        self._add_text_box(slide, title, 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")

        headers = [str(h) for h in (content.get("headers") or [])]
        rows = content.get("rows") or []
        norm_rows = [[str(c) for c in (r or [])] for r in rows]
        n_cols = max([len(headers)] + [len(r) for r in norm_rows] + [2])
        # Pad header + rows to a uniform column count.
        headers = (headers + [""] * n_cols)[:n_cols]
        norm_rows = [(r + [""] * n_cols)[:n_cols] for r in norm_rows]
        n_rows = len(norm_rows) + 1  # +1 header

        left, top, width = 0.6, 1.55, 12.13
        # Cap total height; PowerPoint distributes evenly.
        height = min(5.4, 0.55 + 0.55 * len(norm_rows))

        gfx = slide.shapes.add_table(n_rows, n_cols,
                                     Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        table = gfx.table
        # Disable PowerPoint's default banded styling so our explicit fills show.
        table.first_row = False
        table.horz_banding = False

        # First column wider (the feature/label column).
        first_col_w = Emu(int(Inches(width) * 0.30))
        other_w = Emu(int((Inches(width) - first_col_w) / max(n_cols - 1, 1)))
        table.columns[0].width = first_col_w
        for c in range(1, n_cols):
            table.columns[c].width = other_w

        accent_rgb = self._rgb("accent")
        header_text_hex = self._contrast_on(self.theme.get("accent", "#000000"))
        header_text_rgb = RGBColor.from_string(header_text_hex.lstrip("#"))
        heading_rgb = self._rgb("heading")
        text_rgb = self._rgb("text")
        bg_rgb = self._rgb("bg")
        secondary_rgb = self._rgb("secondary")

        def _fill(cell, rgb):
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb

        def _set(cell, text, color_rgb, bold=False, align=_AL.LEFT, size=13):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Emu(64000)
            cell.margin_right = Emu(64000)
            cell.margin_top = Emu(20000)
            cell.margin_bottom = Emu(20000)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            for r in list(p.runs):
                r.text = ""
            run = p.add_run()
            run.text = text
            run.font.size = _Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color_rgb

        # Header row
        for c in range(n_cols):
            cell = table.cell(0, c)
            _fill(cell, accent_rgb)
            _set(cell, headers[c], header_text_rgb, bold=True,
                 align=_AL.LEFT if c == 0 else _AL.CENTER, size=14)

        # Body rows — zebra: even rows bg, odd rows secondary.
        for ri, row in enumerate(norm_rows):
            for c in range(n_cols):
                cell = table.cell(ri + 1, c)
                _fill(cell, secondary_rgb if ri % 2 == 1 else bg_rgb)
                is_first = c == 0
                _set(cell, row[c],
                     heading_rgb if is_first else text_rgb,
                     bold=is_first,
                     align=_AL.LEFT if is_first else _AL.CENTER,
                     size=12)

    def _render_section_header(self, slide, slide_data):
        """Big section divider — keeps the theme bg, uses accent only as a strip."""
        content = slide_data.get("content", {})
        # Wide left accent stripe (rest of the slide stays on the theme bg).
        self._add_accent_bar(slide, 0.0, 0.0, 0.5, 7.5, color_key="accent")
        if content.get("eyebrow"):
            self._add_text_box(slide, content["eyebrow"].upper(), 1.0, 2.6, 11.33, 0.5,
                               font_size=18, bold=True, color_key="accent")
        title = content.get("title", slide_data.get("title", ""))
        self._add_text_box(slide, title, 1.0, 3.2, 11.33, 1.6,
                           font_size=60, bold=True, color_key="heading")
        # Underline accent
        self._add_accent_bar(slide, 1.0, 4.95, 1.0, 0.08, color_key="accent")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 1.0, 5.2, 11.33, 0.9,
                               font_size=20, color_key="muted")

    def _render_big_number(self, slide, slide_data):
        """Single hero stat — for impact moments."""
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=24, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        self._add_text_box(slide, content.get("value", ""),
                           0.6, 1.8, 12.33, 3.5,
                           font_size=200, bold=True, align=PP_ALIGN.CENTER, color_key="accent")
        self._add_text_box(slide, content.get("label", ""), 0.6, 5.4, 12.33, 0.7,
                           font_size=28, bold=True, align=PP_ALIGN.CENTER, color_key="heading")
        if content.get("context"):
            self._add_text_box(slide, content["context"], 1.5, 6.2, 10.33, 0.7,
                               font_size=15, align=PP_ALIGN.CENTER, color_key="muted")

    def _render_icon_grid(self, slide, slide_data):
        """3x2 or 2x2 cards with icon/emoji + heading + description."""
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        items = content.get("items", [])[:6]
        n = len(items)
        if n == 0:
            return
        cols = 3 if n > 4 else 2
        rows = (n + cols - 1) // cols
        gap = 0.25
        avail_w = 12.33 - gap * (cols - 1)
        avail_h = 5.5 - gap * (rows - 1)
        cw = avail_w / cols
        ch = avail_h / rows
        for i, item in enumerate(items):
            r, c = divmod(i, cols)
            left = 0.5 + c * (cw + gap)
            top = 1.5 + r * (ch + gap)
            self._add_card(slide, left, top, cw, ch, color_key="secondary")
            # Match the browser: small rounded pill (theme.accent at ~15%
            # opacity) sitting behind the icon glyph. PowerPoint shapes can't
            # alpha-blend reliably, so we approximate with a `card`-coloured
            # pill that reads as a subtle backdrop on any theme.
            pill_size = 0.55
            pill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.2), Inches(top + 0.2),
                Inches(pill_size), Inches(pill_size),
            )
            pill.line.fill.background()
            pill.fill.solid()
            pill.fill.fore_color.rgb = self._rgb("card")
            self._add_text_box(slide, item.get("icon", "●"),
                               left + 0.2, top + 0.22,
                               pill_size, pill_size,
                               font_size=22, bold=True,
                               align=PP_ALIGN.CENTER, color_key="accent")
            self._add_text_box(slide, item.get("heading", ""), left + 0.2, top + 0.9,
                               cw - 0.4, 0.5, font_size=16, bold=True, color_key="heading")
            self._add_text_box(slide, item.get("description", ""), left + 0.2, top + 1.45,
                               cw - 0.4, ch - 1.55, font_size=12, color_key="muted")

    def _render_arrow_columns(self, slide, slide_data):
        """3-4 columns each with an arrow icon (→ glyph), heading, description."""
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        items = content.get("items", [])[:4]
        n = max(len(items), 1)
        gap = 0.4
        col_w = (12.33 - gap * (n - 1)) / n
        for i, item in enumerate(items):
            left = 0.5 + i * (col_w + gap)
            # → arrow glyph
            self._add_text_box(slide, "→", left, 1.7, 1.0, 0.8,
                               font_size=36, bold=True, color_key="accent")
            self._add_text_box(slide, item.get("heading", ""),
                               left, 2.7, col_w, 0.8,
                               font_size=22, bold=True, color_key="heading")
            self._add_text_box(slide, item.get("description", ""),
                               left, 3.7, col_w, 3.0,
                               font_size=14, color_key="text")

    def _render_image_with_cards(self, slide, slide_data):
        """Image on left half, eyebrow + title + 3-4 cards on right half."""
        content = slide_data.get("content", {})
        # Left image
        if content.get("image_url"):
            self._add_image(slide, content["image_url"], 0.0, 0.0, 6.5, 7.5)
        else:
            self._add_card(slide, 0.0, 0.0, 6.5, 7.5, color_key="secondary", radius=False)
        # Right column
        eyebrow_txt = content.get("eyebrow", "")
        if eyebrow_txt:
            # Match the browser: rounded chip outlined in accent, text in
            # accent. Sized to the content (heuristic 0.12" per character +
            # padding) so we don't draw an empty box across the column.
            chip_w = min(max(0.6, len(eyebrow_txt) * 0.13), 4.5)
            chip_h = 0.35
            chip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), Inches(0.55), Inches(chip_w), Inches(chip_h),
            )
            chip.fill.background()
            chip.line.color.rgb = self._rgb("accent")
            chip.line.width = Pt(1)
            self._add_text_box(slide, eyebrow_txt.upper(),
                               6.8, 0.58, chip_w, chip_h - 0.05,
                               font_size=11, bold=True,
                               align=PP_ALIGN.CENTER, color_key="accent")
        self._add_text_box(slide, content.get("title", slide_data.get("title", "")),
                           6.8, 1.0, 6.3, 1.8,
                           font_size=32, bold=True, color_key="heading")
        # Cards in 2-col grid
        cards = content.get("cards", [])[:4]
        n = len(cards)
        if n == 0:
            return
        cols = 2
        rows = (n + 1) // 2
        gap = 0.2
        cw = (6.3 - gap) / cols
        ch_total = 4.0
        ch = (ch_total - gap * (rows - 1)) / rows
        for i, c in enumerate(cards):
            r, col = divmod(i, cols)
            # Last odd card spans 2 cols
            if i == n - 1 and n % 2 == 1:
                left = 6.8
                width = 6.3
            else:
                left = 6.8 + col * (cw + gap)
                width = cw
            top = 3.0 + r * (ch + gap)
            self._add_card(slide, left, top, width, ch, color_key="card")
            self._add_text_box(slide, c.get("heading", ""),
                               left + 0.2, top + 0.2, width - 0.4, 0.6,
                               font_size=16, bold=True, color_key="heading")
            self._add_text_box(slide, c.get("description", ""),
                               left + 0.2, top + 0.9, width - 0.4, ch - 1.0,
                               font_size=11, color_key="text")

    def _render_team_image_grid(self, slide, slide_data):
        """3 portrait cards each with image + name + bio."""
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", "Team"),
                           0.6, 0.4, 12.33, 0.9,
                           font_size=38, bold=True, color_key="heading")
        members = content.get("members", [])[:3]
        n = max(len(members), 1)
        gap = 0.3
        col_w = (12.33 - gap * (n - 1)) / n
        img_h = 3.0
        for i, m in enumerate(members):
            left = 0.5 + i * (col_w + gap)
            top = 1.6
            if m.get("image_url"):
                self._add_image(slide, m["image_url"], left, top, col_w, img_h)
            else:
                self._add_card(slide, left, top, col_w, img_h, color_key="secondary")
                self._add_text_box(slide,
                                   (m.get("name", "?") or "?")[0].upper(),
                                   left, top + 0.8, col_w, 1.4,
                                   font_size=64, bold=True,
                                   align=PP_ALIGN.CENTER, color_key="accent")
            self._add_text_box(slide, m.get("name", ""),
                               left, top + img_h + 0.2, col_w, 0.6,
                               font_size=20, bold=True, color_key="heading")
            self._add_text_box(slide, m.get("bio", ""),
                               left, top + img_h + 0.9, col_w, 2.4,
                               font_size=12, color_key="text")

    def _render_process_steps(self, slide, slide_data):
        """Numbered horizontal steps with arrow separators."""
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        steps = content.get("steps", [])[:5]
        n = max(len(steps), 1)
        arrow_w = 0.4
        gap_total = arrow_w * (n - 1) + 0.2 * (n - 1)
        cw = (12.33 - gap_total) / n
        x = 0.5
        for i, step in enumerate(steps):
            top = 2.0
            self._add_card(slide, x, top, cw, 4.5, color_key="secondary")
            # Number bubble
            num_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(top + 0.25),
                Inches(0.65), Inches(0.65)
            )
            num_shape.line.fill.background()
            num_shape.fill.solid()
            num_shape.fill.fore_color.rgb = self._rgb("accent")
            self._add_text_box(slide, str(i + 1).zfill(2),
                               x + 0.25, top + 0.30, 0.65, 0.55,
                               font_size=14, bold=True,
                               align=PP_ALIGN.CENTER, color_key="bg")
            self._add_text_box(slide, step.get("heading", ""),
                               x + 0.25, top + 1.1, cw - 0.5, 0.7,
                               font_size=16, bold=True, color_key="heading")
            self._add_text_box(slide, step.get("description", ""),
                               x + 0.25, top + 1.9, cw - 0.5, 2.4,
                               font_size=11, color_key="text")
            x += cw + 0.2
            if i < n - 1:
                # Arrow glyph between cards
                self._add_text_box(slide, "→", x - 0.1, top + 1.8, arrow_w, 0.8,
                                   font_size=24, bold=True,
                                   align=PP_ALIGN.CENTER, color_key="accent")
                x += arrow_w

    def _render_pyramid(self, slide, slide_data):
        """Stacked bars forming a pyramid (top = narrowest). Each level shows
        a bold label on the left and a description on the right so the
        exported deck matches the on-screen pyramid renderer."""
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        levels = content.get("levels", [])[:4]
        n = max(len(levels), 1)
        # Top is narrowest, bottom is widest
        top_w_pct = 0.30
        bot_w_pct = 0.90
        slide_w = 12.33
        row_h = 0.95
        gap = 0.15
        avail_h = 5.6
        total_h = row_h * n + gap * (n - 1)
        start_top = 1.8 + (avail_h - total_h) / 2
        for i, lv in enumerate(levels):
            pct = top_w_pct + (i / max(n - 1, 1)) * (bot_w_pct - top_w_pct)
            w = slide_w * pct
            left = 0.5 + (slide_w - w) / 2
            top = start_top + i * (row_h + gap)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(w), Inches(row_h)
            )
            shape.line.fill.background()
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb("accent")

            label = lv.get("label", "") or ""
            description = lv.get("description", "") or ""
            # Match the on-screen renderer: label-on-left, description-on-right,
            # both in the slide's bg colour (high-contrast against accent).
            if description:
                lbl_w = max(w * 0.4, 1.4)
                desc_w = w - lbl_w - 0.4
                self._add_text_box(slide, label,
                                   left + 0.25, top + 0.18, lbl_w, row_h - 0.3,
                                   font_size=14, bold=True,
                                   align=PP_ALIGN.LEFT, color_key="bg")
                self._add_text_box(slide, description,
                                   left + 0.25 + lbl_w, top + 0.12,
                                   desc_w, row_h - 0.2,
                                   font_size=10, bold=False,
                                   align=PP_ALIGN.RIGHT, color_key="bg")
            else:
                self._add_text_box(slide, label,
                                   left + 0.3, top + 0.18, w - 0.6, row_h - 0.3,
                                   font_size=15, bold=True,
                                   align=PP_ALIGN.CENTER, color_key="bg")

    def _render_cta(self, slide, slide_data):
        """Call-to-action / closing slide."""
        content = slide_data.get("content", {})
        self._add_card(slide, 0.0, 0.0, 13.33, 7.5, color_key="bg", radius=False)
        self._add_text_box(slide, content.get("title", ""), 1.0, 2.0, 11.33, 1.5,
                           font_size=54, bold=True, align=PP_ALIGN.CENTER, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 1.5, 3.5, 10.33, 1.0,
                               font_size=22, align=PP_ALIGN.CENTER, color_key="muted")
        # Button shape
        if content.get("button_label"):
            btn = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.16), Inches(5.0), Inches(3.0), Inches(0.8)
            )
            btn.line.fill.background()
            btn.fill.solid()
            btn.fill.fore_color.rgb = self._rgb("accent")
            # Use heading color (always high contrast vs accent) instead of
            # bg, which can be ~the same brightness as accent in some themes.
            self._add_text_box(slide, content["button_label"], 5.16, 5.1, 3.0, 0.6,
                               font_size=20, bold=True, align=PP_ALIGN.CENTER, color_key="heading")
        if content.get("contact"):
            self._add_text_box(slide, content["contact"], 1.0, 6.4, 11.33, 0.5,
                               font_size=14, align=PP_ALIGN.CENTER, color_key="muted")

    # ── Helpers shared by diagrams ─────────────────────────────────────────
    def _add_filled_shape(self, slide, shape_type, left, top, width, height,
                          fill_key="accent", fill_alpha=1.0, line=False):
        """Filled rectangle/oval/triangle helper used by diagrams. Alpha is
        emulated via colour blend (PPT shape opacity is a beast; blending
        toward bg gives the same visual)."""
        from pptx.dml.color import RGBColor
        shape = slide.shapes.add_shape(
            shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if line:
            shape.line.color.rgb = self._rgb("bg")
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        shape.fill.solid()
        if fill_alpha >= 0.99:
            shape.fill.fore_color.rgb = self._rgb(fill_key)
        else:
            base = self.theme.get(fill_key, "#000000").lstrip("#")
            bg = self.theme.get("bg", "#FFFFFF").lstrip("#")
            br, bgg, bb = int(base[0:2], 16), int(base[2:4], 16), int(base[4:6], 16)
            wr, wg, wb = int(bg[0:2], 16), int(bg[2:4], 16), int(bg[4:6], 16)
            r = int(wr + (br - wr) * fill_alpha)
            g = int(wg + (bgg - wg) * fill_alpha)
            b = int(wb + (bb - wb) * fill_alpha)
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
        return shape

    # ── 1. Funnel ─────────────────────────────────────────────────────────
    def _render_funnel(self, slide, slide_data):
        """Real-trapezoid funnel that matches the SVG renderer's polygon
        slices: each stage is a freeform 4-point polygon whose top width
        equals the previous stage's bottom width (so the trapezoids tile
        seamlessly into a continuous funnel)."""
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        stages = (content.get("stages") or [])[:5]
        n = max(len(stages), 1)
        TOP_W = 11.0
        BOT_W = 4.0
        H_AVAIL = 5.4
        ROW_H = H_AVAIL / n
        TOP = 1.6
        cx = 13.33 / 2
        # One legible text colour for every slice (computed against accent).
        txt_hex = self._contrast_on(self.theme.get("accent", "#000000"))
        for i, stage in enumerate(stages):
            t0 = i / n
            t1 = (i + 1) / n
            w0 = TOP_W - (TOP_W - BOT_W) * t0
            w1 = TOP_W - (TOP_W - BOT_W) * t1
            y0 = TOP + i * ROW_H
            y1 = y0 + ROW_H - 0.05
            alpha = 0.55 + (i / max(n - 1, 1)) * 0.45
            # Draw the trapezoid as a freeform polygon: 4 points, clockwise.
            self._add_polygon(
                slide,
                points=[
                    (cx - w0 / 2, y0),
                    (cx + w0 / 2, y0),
                    (cx + w1 / 2, y1),
                    (cx - w1 / 2, y1),
                ],
                fill_key="accent",
                fill_alpha=alpha,
            )
            # Constrain text to the NARROW (bottom) edge so it never spills
            # past the trapezoid. Text wraps inside the box.
            text_w = max(w1 - 0.5, 1.6)
            desc = str(stage.get("description", "") or "")
            # Label sits in the upper half, description in the lower half.
            self._add_text_box(slide, str(stage.get("label", "")),
                               cx - text_w / 2, y0 + ROW_H * 0.12,
                               text_w, ROW_H * 0.42,
                               font_size=15, bold=True,
                               align=PP_ALIGN.CENTER, color_hex=txt_hex)
            if desc:
                self._add_text_box(slide, desc,
                                   cx - text_w / 2, y0 + ROW_H * 0.52,
                                   text_w, ROW_H * 0.4,
                                   font_size=10,
                                   align=PP_ALIGN.CENTER, color_hex=txt_hex)

    def _add_polygon(self, slide, points, fill_key="accent", fill_alpha=1.0):
        """Add a freeform polygon by tracing the supplied (x, y) inch points.
        Uses python-pptx's FreeformBuilder so the result is a vector shape
        editable inside PowerPoint."""
        from pptx.dml.color import RGBColor
        # build_freeform expects the start point in LOCAL units and `scale`
        # to convert local→EMU. Passing Inches() (already EMU) here would
        # double-convert and fling the shape millions of inches off-slide
        # (this was why the funnel rendered invisible). Pass raw inch floats
        # with scale = EMU-per-inch instead.
        EMU_PER_IN = 914400
        first_x, first_y = points[0]
        builder = slide.shapes.build_freeform(
            float(first_x), float(first_y), scale=EMU_PER_IN,
        )
        # add_line_segments takes ABSOLUTE local (inch) coordinates.
        builder.add_line_segments(
            [(float(x), float(y)) for x, y in points[1:]], close=True
        )
        shape = builder.convert_to_shape()
        shape.line.fill.background()
        shape.fill.solid()
        if fill_alpha >= 0.99:
            shape.fill.fore_color.rgb = self._rgb(fill_key)
        else:
            base = self.theme.get(fill_key, "#000000").lstrip("#")
            bgh = self.theme.get("bg", "#FFFFFF").lstrip("#")
            br, bgg, bb = int(base[0:2], 16), int(base[2:4], 16), int(base[4:6], 16)
            wr, wg, wb = int(bgh[0:2], 16), int(bgh[2:4], 16), int(bgh[4:6], 16)
            r = int(wr + (br - wr) * fill_alpha)
            g = int(wg + (bgg - wg) * fill_alpha)
            b = int(wb + (bb - wb) * fill_alpha)
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
        return shape

    # ── 2. Concentric circles ─────────────────────────────────────────────
    def _render_concentric_circles(self, slide, slide_data):
        from pptx.enum.shapes import MSO_SHAPE
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        layers = (content.get("layers") or [])[:5]
        n = max(len(layers), 1)
        cx, cy = 4.5, 4.2
        R = 2.5
        for i in range(n):
            r = R - (i * R) / n
            alpha = 0.25 + (i / max(n - 1, 1)) * 0.7
            self._add_filled_shape(
                slide, MSO_SHAPE.OVAL,
                cx - r, cy - r, r * 2, r * 2,
                fill_key="accent", fill_alpha=alpha,
            )
        # Right-hand legend
        for i, lv in enumerate(layers):
            y = 1.8 + i * (5.0 / max(n, 1))
            self._add_filled_shape(
                slide, MSO_SHAPE.OVAL, 7.6, y, 0.25, 0.25,
                fill_key="accent",
                fill_alpha=0.3 + (i / max(n - 1, 1)) * 0.7,
            )
            self._add_text_box(slide, str(lv.get("label", "")),
                               7.95, y - 0.05, 5.0, 0.4,
                               font_size=15, bold=True, color_key="heading")
            if lv.get("description"):
                self._add_text_box(slide, str(lv["description"]),
                                   7.95, y + 0.3, 5.0, 0.5,
                                   font_size=10, color_key="muted")

    # ── 3. Venn (2-set) ───────────────────────────────────────────────────
    def _render_venn(self, slide, slide_data):
        from pptx.enum.shapes import MSO_SHAPE
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        set_a = content.get("set_a", {}) or {}
        set_b = content.get("set_b", {}) or {}
        overlap_label = str(content.get("overlap_label", "") or "")
        overlap_items = content.get("overlap_items") or []
        r = 2.2
        cy = 4.2
        cxa = 4.6
        cxb = 7.0
        self._add_filled_shape(slide, MSO_SHAPE.OVAL,
                               cxa - r, cy - r, r * 2, r * 2,
                               fill_key="accent", fill_alpha=0.4)
        self._add_filled_shape(slide, MSO_SHAPE.OVAL,
                               cxb - r, cy - r, r * 2, r * 2,
                               fill_key="accent", fill_alpha=0.4)
        # Set labels above
        self._add_text_box(slide, str(set_a.get("label", "")),
                           cxa - r, cy - r - 0.55, r * 2, 0.45,
                           font_size=18, bold=True,
                           align=PP_ALIGN.CENTER, color_key="heading")
        self._add_text_box(slide, str(set_b.get("label", "")),
                           cxb - r, cy - r - 0.55, r * 2, 0.45,
                           font_size=18, bold=True,
                           align=PP_ALIGN.CENTER, color_key="heading")
        # Items in each lobe
        for i, it in enumerate((set_a.get("items") or [])[:4]):
            self._add_text_box(slide, f"• {it}",
                               cxa - r - 0.1, cy - 0.6 + i * 0.4, 1.7, 0.4,
                               font_size=12, bold=True,
                               align=PP_ALIGN.RIGHT, color_key="bg")
        for i, it in enumerate((set_b.get("items") or [])[:4]):
            self._add_text_box(slide, f"• {it}",
                               cxb + r - 1.6, cy - 0.6 + i * 0.4, 1.7, 0.4,
                               font_size=12, bold=True,
                               align=PP_ALIGN.LEFT, color_key="bg")
        # Overlap
        if overlap_label:
            self._add_text_box(slide, overlap_label,
                               (cxa + cxb) / 2 - 1.2, cy - 0.5, 2.4, 0.5,
                               font_size=14, bold=True,
                               align=PP_ALIGN.CENTER, color_key="bg")
        for i, it in enumerate(overlap_items[:3]):
            self._add_text_box(slide, f"• {it}",
                               (cxa + cxb) / 2 - 1.2, cy + 0.1 + i * 0.35, 2.4, 0.35,
                               font_size=11, align=PP_ALIGN.CENTER, color_key="bg")

    # ── 4. Target / bullseye ──────────────────────────────────────────────
    def _render_target(self, slide, slide_data):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Pt
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        rings = (content.get("rings") or [])[:5]
        n = max(len(rings), 1)
        cx, cy = 3.8, 4.2
        R = 2.4
        for i in range(n):
            r = R * (1 - i / n)
            alpha = 0.3 + (i / max(n - 1, 1)) * 0.65
            self._add_filled_shape(
                slide, MSO_SHAPE.OVAL,
                cx - r, cy - r, r * 2, r * 2,
                fill_key="accent", fill_alpha=alpha, line=True,
            )
        # Crosshair lines through the centre — matches the SVG renderer.
        from pptx.dml.color import RGBColor
        muted_hex = self.theme.get("muted", "#94A3B8").lstrip("#")
        cross_color = RGBColor.from_string(muted_hex)
        h = slide.shapes.add_connector(1, Inches(cx - R - 0.2), Inches(cy),
                                       Inches(cx + R + 0.2), Inches(cy))
        h.line.color.rgb = cross_color
        h.line.width = Pt(0.75)
        v = slide.shapes.add_connector(1, Inches(cx), Inches(cy - R - 0.2),
                                       Inches(cx), Inches(cy + R + 0.2))
        v.line.color.rgb = cross_color
        v.line.width = Pt(0.75)
        # Legend right
        for i, ring in enumerate(rings):
            y = 1.8 + i * (5.0 / max(n, 1))
            self._add_filled_shape(
                slide, MSO_SHAPE.OVAL, 7.4, y, 0.22, 0.22,
                fill_key="accent",
                fill_alpha=0.3 + (i / max(n - 1, 1)) * 0.7,
            )
            self._add_text_box(slide, str(ring.get("label", "")),
                               7.7, y - 0.07, 5.4, 0.4,
                               font_size=15, bold=True, color_key="heading")
            if ring.get("description"):
                self._add_text_box(slide, str(ring["description"]),
                                   7.7, y + 0.28, 5.4, 0.5,
                                   font_size=10, color_key="muted")

    # ── 5. Connected circles ──────────────────────────────────────────────
    def _render_connected_circles(self, slide, slide_data):
        from pptx.enum.shapes import MSO_SHAPE
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.8,
                           font_size=30, bold=True, color_key="heading")
        self._add_accent_bar(slide, 0.6, 1.2, 0.7, 0.06, color_key="accent")
        nodes = (content.get("nodes") or [])[:5]
        n = max(len(nodes), 1)
        margin = 1.0
        circ = 1.3
        slide_w = 13.33
        span = (slide_w - margin * 2 - circ) / max(n - 1, 1) if n > 1 else 0
        cy = 3.5
        # Connector — accent stripe behind the circles
        if n > 1:
            self._add_accent_bar(slide,
                                 margin + circ / 2, cy + circ / 2 - 0.03,
                                 span * (n - 1), 0.06,
                                 color_key="accent")
        for i, node in enumerate(nodes):
            x = margin + i * span
            alpha = 0.55 + (i / max(n - 1, 1)) * 0.45
            self._add_filled_shape(
                slide, MSO_SHAPE.OVAL, x, cy, circ, circ,
                fill_key="accent", fill_alpha=alpha,
            )
            # Center the label vertically inside the circle (the SVG renderer
            # uses textAnchor middle + baseline-aware y; mimic with a box
            # that's the height of the circle and centered alignment).
            label_h = 0.45
            self._add_text_box(slide, str(node.get("label", "")),
                               x, cy + (circ - label_h) / 2,
                               circ, label_h,
                               font_size=14, bold=True,
                               align=PP_ALIGN.CENTER, color_key="bg")
            if node.get("description"):
                self._add_text_box(slide, str(node["description"]),
                                   x - 0.5, cy + circ + 0.15, circ + 1.0, 1.6,
                                   font_size=11,
                                   align=PP_ALIGN.CENTER, color_key="text")

    # ── Charts (native python-pptx CategoryChartData / XyChartData) ───────
    @staticmethod
    def _num_list(raw) -> list:
        """Coerce chart series data into a list of floats. Accepts a real
        list of numbers, OR a comma/space-joined string like "12, 8, 20"
        (which older slides stored before the flatten fix). Non-numeric
        tokens become 0.0 so the chart still draws."""
        def _to_float(x):
            if isinstance(x, (int, float)):
                return float(x)
            try:
                return float(str(x).strip())
            except (TypeError, ValueError):
                return 0.0
        if isinstance(raw, list):
            return [_to_float(v) for v in raw]
        if isinstance(raw, (int, float)):
            # Single scalar — used by pie/donut where each slice value is one
            # number (not an array).
            return [float(raw)]
        if isinstance(raw, str) and raw.strip():
            import re as _re
            return [_to_float(tok) for tok in _re.split(r"[,\s]+", raw.strip()) if tok]
        return []

    def _theme_chart(self, chart, axes=True):
        """Make a native chart match the browser SVG look on any theme:
        hide the auto chart title, recolour every text element to the theme's
        text/muted colour (PowerPoint defaults to near-black, invisible on
        dark themes), and strip the plot/chart border boxes the SVG doesn't
        have."""
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn

        text_rgb = RGBColor.from_string(self.theme.get("text", "#1A1A1A").lstrip("#"))
        muted_rgb = RGBColor.from_string(self.theme.get("muted", "#6B7280").lstrip("#"))

        chart.has_title = False

        # Default chart font colour (legend etc.)
        try:
            chart.font.size = Pt(11)
            chart.font.color.rgb = text_rgb
        except Exception:
            pass

        if axes:
            for ax, col in ((getattr(chart, "category_axis", None), muted_rgb),
                            (getattr(chart, "value_axis", None), muted_rgb)):
                if ax is None:
                    continue
                try:
                    ax.tick_labels.font.color.rgb = col
                    ax.tick_labels.font.size = Pt(11)
                except Exception:
                    pass
                # Soften gridlines / axis lines toward the border colour.
                try:
                    ax.format.line.color.rgb = RGBColor.from_string(
                        self.theme.get("border", "#334155").lstrip("#"))
                except Exception:
                    pass

        # Remove the chartSpace outer border (the box around the whole chart)
        # by injecting an empty line (<a:ln><a:noFill/></a:ln>) into spPr.
        try:
            chartSpace = chart._chartSpace
            spPr = chartSpace.find(qn("c:spPr"))
            if spPr is None:
                spPr = chartSpace.makeelement(qn("c:spPr"), {})
                chartSpace.insert(0, spPr)
            # Clear any existing line, add noFill line.
            for existing in spPr.findall(qn("a:ln")):
                spPr.remove(existing)
            ln = spPr.makeelement(qn("a:ln"), {})
            noFill = ln.makeelement(qn("a:noFill"), {})
            ln.append(noFill)
            spPr.append(ln)
        except Exception:
            pass

    def _add_native_chart(self, slide, chart_type, content, left, top, width, height):
        """Build a CategoryChartData → add as native chart object so PowerPoint
        treats it as editable. Both bar/line/area share the categorical
        structure; pie/donut go through a separate path with one series."""
        from pptx.chart.data import CategoryChartData
        from pptx.util import Inches as I
        from pptx.dml.color import RGBColor

        categories = [str(c) for c in (content.get("categories") or [])]
        if not categories:
            categories = ["—"]
        series = content.get("series") or []
        if not series:
            series = [{"name": "Series 1", "data": [0] * len(categories)}]

        data = CategoryChartData()
        data.categories = categories
        palette = ["DC2626", "10B981", "F59E0B", "8B5CF6", "EC4899", "06B6D4"]
        for idx, s in enumerate(series):
            name = str(s.get("name", f"Series {idx+1}"))
            values = self._num_list(s.get("data"))
            # Pad/truncate to match category count so python-pptx is happy.
            if len(values) < len(categories):
                values += [0.0] * (len(categories) - len(values))
            elif len(values) > len(categories):
                values = values[:len(categories)]
            data.add_series(name, values)

        chart_shape = slide.shapes.add_chart(
            chart_type, I(left), I(top), I(width), I(height), data
        )
        chart = chart_shape.chart
        # Theme the chart: paint each series with our palette.
        for i, plot_series in enumerate(chart.plots[0].series):
            color = (series[i].get("color") if i < len(series) else None) or palette[i % len(palette)]
            color = color.lstrip("#")
            fmt = plot_series.format
            fmt.fill.solid()
            fmt.fill.fore_color.rgb = RGBColor.from_string(color)
            try:
                fmt.line.color.rgb = RGBColor.from_string(color)
                fmt.line.width = Pt(2.5)
            except Exception:
                pass
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            try:
                from pptx.enum.chart import XL_LEGEND_POSITION
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            except Exception:
                pass
        self._theme_chart(chart, axes=True)
        return chart

    def _add_pie_chart(self, slide, content, left, top, width, height, donut=False):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches as I
        from pptx.dml.color import RGBColor

        slices = content.get("slices") or []
        if not slices:
            slices = [{"label": "—", "value": 1}]
        labels = [str(s.get("label", "")) for s in slices]
        values = [self._num_list(s.get("value"))[0] if self._num_list(s.get("value")) else 0.0
                  for s in slices]

        data = CategoryChartData()
        data.categories = labels
        data.add_series("Series", values)

        chart_type = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE
        chart_shape = slide.shapes.add_chart(
            chart_type, I(left), I(top), I(width), I(height), data
        )
        chart = chart_shape.chart
        palette = ["DC2626", "10B981", "F59E0B", "8B5CF6", "EC4899", "06B6D4"]
        for i, pt in enumerate(chart.plots[0].series[0].points):
            color = (slices[i].get("color") if i < len(slices) else None) or palette[i % len(palette)]
            color = str(color).lstrip("#")
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RGBColor.from_string(color)
        chart.has_legend = True
        try:
            from pptx.enum.chart import XL_LEGEND_POSITION
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
        except Exception:
            pass
        # Pie/donut have no axes — only title/legend/border theming.
        self._theme_chart(chart, axes=False)
        return chart

    def _render_bar_chart(self, slide, slide_data):
        from pptx.enum.chart import XL_CHART_TYPE
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 0.6, 1.15, 12.33, 0.4,
                               font_size=14, color_key="muted")
        self._add_accent_bar(slide, 0.6, 1.65, 0.7, 0.06, color_key="accent")
        self._add_native_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED,
                               content, 0.6, 1.9, 12.13, 5.2)

    def _render_line_chart(self, slide, slide_data):
        from pptx.enum.chart import XL_CHART_TYPE
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 0.6, 1.15, 12.33, 0.4,
                               font_size=14, color_key="muted")
        self._add_accent_bar(slide, 0.6, 1.65, 0.7, 0.06, color_key="accent")
        self._add_native_chart(slide, XL_CHART_TYPE.LINE_MARKERS,
                               content, 0.6, 1.9, 12.13, 5.2)

    def _render_area_chart(self, slide, slide_data):
        from pptx.enum.chart import XL_CHART_TYPE
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 0.6, 1.15, 12.33, 0.4,
                               font_size=14, color_key="muted")
        self._add_accent_bar(slide, 0.6, 1.65, 0.7, 0.06, color_key="accent")
        self._add_native_chart(slide, XL_CHART_TYPE.AREA,
                               content, 0.6, 1.9, 12.13, 5.2)

    def _render_pie_chart(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 0.6, 1.15, 12.33, 0.4,
                               font_size=14, color_key="muted")
        self._add_accent_bar(slide, 0.6, 1.65, 0.7, 0.06, color_key="accent")
        self._add_pie_chart(slide, content, 0.6, 1.9, 12.13, 5.2, donut=False)

    def _render_donut_chart(self, slide, slide_data):
        content = slide_data.get("content", {})
        self._add_text_box(slide, content.get("title", ""), 0.6, 0.4, 12.33, 0.7,
                           font_size=30, bold=True, color_key="heading")
        if content.get("subtitle"):
            self._add_text_box(slide, content["subtitle"], 0.6, 1.15, 12.33, 0.4,
                               font_size=14, color_key="muted")
        self._add_accent_bar(slide, 0.6, 1.65, 0.7, 0.06, color_key="accent")
        self._add_pie_chart(slide, content, 0.6, 1.9, 12.13, 5.2, donut=True)
